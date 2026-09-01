"""Сборка .pptx в стиле минимализм Ч/Б."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.ollama_client import Slide
from app.utils.logging import logger

# Палитра: минимализм Ч/Б
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)   # белый
COLOR_FG = RGBColor(0x10, 0x10, 0x10)   # почти чёрный
COLOR_MUTED = RGBColor(0x55, 0x55, 0x55)
COLOR_ACCENT = RGBColor(0xC8, 0x1D, 0x1D)  # красный акцент

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def build_pptx(
    slides: list[Slide],
    *,
    title: str,
    photo_paths: Iterable[str | Path] = (),
    output_path: str | Path,
) -> Path:
    """Собирает .pptx и возвращает путь к файлу."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    photos = list(photo_paths)
    n = len(slides)

    # Титульный
    _add_title_slide(prs, title=title, subtitle=slides[0].title if slides else "")

    # Основные слайды
    photo_plan = _distribute_photos(n, len(photos))
    for idx, slide in enumerate(slides):
        photo = photos[photo_plan[idx]] if photo_plan[idx] is not None else None
        _add_content_slide(prs, slide, photo)

    prs.save(output_path)
    logger.info("PPTX built: {} slides, {} photos → {}", n, len(photos), output_path)
    return output_path


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, COLOR_BG)

    # Акцентная полоса
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(4.7),
        SLIDE_W, Inches(0.08),
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.name = "Calibri"; run.font.size = Pt(40); run.font.bold = True
    run.font.color.rgb = COLOR_FG

    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = "Calibri"; r2.font.size = Pt(20)
        r2.font.color.rgb = COLOR_MUTED


def _process_image(photo_path: str | Path, max_w: float, max_h: float) -> tuple[BytesIO, float, float]:
    """Ресайзит изображение с сохранением пропорций и возвращает (поток, итоговая_ширина, итоговая_высота).

    max_w и max_h передаются в дюймах.
    """
    with Image.open(photo_path) as img:
        # Конвертируем в RGB, если нужно (для сохранения в JPG/PNG)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        # Используем thumbnail для сохранения пропорций.
        # Чтобы thumbnail работал, нам нужны значения в пикселях.
        # Принимаем стандарт 96 DPI для расчета максимального размера.
        max_px_w = int(max_w * 96)
        max_px_h = int(max_h * 96)

        img.thumbnail((max_px_w, max_px_h))

        # Вычисляем итоговые размеры в дюймах для pptx
        curr_w, curr_h = img.size
        final_w = Inches(curr_w / 96)
        final_h = Inches(curr_h / 96)

        bio = BytesIO()
        img.save(bio, format="JPEG", quality=85)
        bio.seek(0)
        return bio, final_w, final_h

def _add_content_slide(prs: Presentation, slide_data: Slide, photo_path: str | Path | None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, COLOR_BG)


    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = slide_data.title
    run.font.name = "Calibri"; run.font.size = Pt(28); run.font.bold = True
    run.font.color.rgb = COLOR_FG

    # Акцентная полоса под заголовком
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.55),
        Inches(1.2), Inches(0.06),
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Разметка: слева буллеты, справа фото (если есть)
    bullets_w = Inches(7.6) if photo_path else Inches(12.1)
    photo_box = None
    if photo_path and Path(photo_path).exists():
        photo_box = (Inches(8.4), Inches(1.9), Inches(4.4), Inches(4.6))

    # Буллеты
    bullet_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), bullets_w, Inches(5.0))
    btf = bullet_box.text_frame; btf.word_wrap = True
    for i, bullet in enumerate(slide_data.bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        run = para.add_run()
        run.text = f"• {bullet}"
        run.font.name = "Calibri"; run.font.size = Pt(20)
        run.font.color.rgb = COLOR_FG

    # Фото
    if photo_box is not None:
        x, y, max_w, max_h = photo_box
        try:
            img_bio, final_w, final_h = _process_image(photo_path, max_w, max_h)
            slide.shapes.add_picture(img_bio, x, y, width=final_w, height=final_h)
        except Exception as exc:  # noqa: BLE001 — fallback чтобы не валить генерацию
            logger.warning("Не удалось вставить фото {}: {}", photo_path, exc)


def _fill_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def _distribute_photos(slides_n: int, photos_n: int) -> list[int | None]:
    """Распределяет индексы фото по слайдам основной части.

    Возвращает список длиной slides_n: индекс фото в ``photos`` или None.
    """
    if photos_n == 0 or slides_n == 0:
        return [None] * slides_n
    if photos_n >= slides_n:
        # На первые photos_n слайдов кладём по фото, остальные None
        return list(range(photos_n)) + [None] * (slides_n - photos_n)
    # photos_n < slides_n: распределяем равномерно по 2..N
    result: list[int | None] = [None] * slides_n
    if slides_n <= 1:
        return [0] if photos_n > 0 else [None] * slides_n
    step = (slides_n - 1) / max(photos_n, 1)
    for i in range(photos_n):
        idx = min(slides_n - 1, int(round(1 + i * step)))
        if result[idx] is None:
            result[idx] = i
        else:
            # Ищем ближайший свободный слот
            for shift in (1, -1, 2, -2):
                cand = idx + shift
                if 0 <= cand < slides_n and result[cand] is None:
                    result[cand] = i
                    break
    return result


def slides_to_text(slides: list[Slide], *, title: str) -> str:
    """Генерирует текстовое представление презентации (.txt)."""
    lines = [title, "=" * len(title), ""]
    for i, s in enumerate(slides, 1):
        lines.append(f"Слайд {i}. {s.title}")
        for b in s.bullets:
            lines.append(f"  • {b}")
        lines.append("")
    return "\n".join(lines)


def make_text_file(slides: list[Slide], *, title: str) -> BytesIO:
    """Возвращает BytesIO с текстом презентации (для отправки в Telegram)."""
    bio = BytesIO(slides_to_text(slides, title=title).encode("utf-8"))
    bio.name = "presentation.txt"  # type: ignore[attr-defined]
    bio.seek(0)
    return bio


__all__ = ["build_pptx", "slides_to_text", "make_text_file"]